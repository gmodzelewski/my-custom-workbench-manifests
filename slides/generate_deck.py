#!/usr/bin/env python3
"""Generate custom-workbench-demo.pptx from slides/deck-outline.md structure."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "custom-workbench-demo.pptx"

ACCENT = RGBColor(0xCC, 0x00, 0x00)
TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x1A)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)

SLIDES: list[tuple[str, list[str]]] = [
    (
        "Custom Workbench Images for OpenShift AI",
        [
            "Jupyter workbench + OpenCode CLI",
            "45-minute session",
        ],
    ),
    (
        "Agenda",
        [
            "OpenShift AI tour (10 min)",
            "Custom images + CI (10 min)",
            "Live demo (18 min)",
            "Wrap-up + Q&A (7 min)",
        ],
    ),
    (
        "Who this is for",
        [
            "Platform engineers registering notebook images",
            "Data scientists who need consistent dev environments",
            "Prerequisites: RHOAI 3.4, cluster admin for ImageStream",
        ],
    ),
    (
        "OpenShift AI in one diagram",
        [
            "DataScienceCluster (operators)",
            "Data Science Projects (tenancy)",
            "Workbenches, model serving, pipelines",
        ],
    ),
    (
        "Data Science Projects",
        [
            "Namespace + dashboard visibility",
            "Shared connections (S3, DB)",
            "Team collaboration boundary",
        ],
    ),
    (
        "Workbenches",
        [
            "JupyterLab and VS Code options",
            "PVC for persistent /opt/app-root/src",
            "Stock images vs bring-your-own",
        ],
    ),
    (
        "Connected resources",
        [
            "PVCs, Secrets, ConfigMaps",
            "Git integration from JupyterLab",
            "Hardware profiles (CPU/GPU)",
        ],
    ),
    (
        "Why customize an image?",
        [
            "Reproducible dependencies across the team",
            "CLI tools, system packages, pinned versions",
            "Same stack in workbench and automation",
        ],
    ),
    (
        "Extension patterns",
        [
            "pip install in notebook (quick, not reproducible)",
            "Extend base image in Containerfile (recommended)",
            "Separate runtime image for pipelines",
        ],
    ),
    (
        "Build pipeline overview",
        [
            "git push → Tekton Pipeline → Quay → ImageStream",
            "Local bootstrap: podman build + push",
            "Dashboard picks up imported tags",
        ],
    ),
    (
        "Dev loop in the workbench",
        [
            "Clone repo in Jupyter Git",
            "Edit Containerfile, commit, push",
            "Stop/start workbench to pick up new image",
        ],
    ),
    (
        "Pipelines: build vs run",
        [
            "Build images: OpenShift Pipelines / Tekton (or podman locally)",
            "Run ML steps: Data Science Pipelines (Kubeflow)",
            "KFP uses pre-built runtime base_image",
        ],
    ),
    (
        "Dashboard registration",
        [
            "Settings → Notebook images",
            "Custom version tag 1.0 (not platform 3.4)",
            "workbench-image-recommended: true",
            "notebook-build-commit must match workbench",
        ],
    ),
    (
        "Demo: what we built",
        [
            "Custom Jupyter Data Science image",
            "OpenCode CLI in terminal",
            "Sample Python app for agent demo",
            "Optional KFP verify pipeline",
        ],
    ),
    (
        "LIVE DEMO",
        [],
    ),
    (
        "OpenCode CLI",
        [
            "opencode run for non-interactive use",
            "Provider auth via env or opencode auth",
            "Agent for code explanation and refactors",
        ],
    ),
    (
        "Operational concerns",
        [
            "Pin base image and tool versions",
            "Scan images for CVEs",
            "Air-gap: vendor OpenCode binary in repo",
            "Rebuild on dependency updates",
        ],
    ),
    (
        "When not to customize",
        [
            "One-off pip install is enough",
            "Short workshops with stock image",
            "Cost of maintaining image > benefit",
        ],
    ),
    (
        "Runtime images for pipelines",
        [
            "Containerfile.runtime without Jupyter",
            "Same packages as workbench where possible",
            "Pipeline component base_image points at runtime",
        ],
    ),
    (
        "Summary",
        [
            "Extend stock RHOAI image with a thin Containerfile",
            "Push to Quay; register via ImageStream",
            "Develop from git; use OpenCode in the terminal",
        ],
    ),
    (
        "Links",
        [
            "Repo: my-custom-workbench",
            "opencode.ai/docs/cli",
            "RHOAI docs: managing notebook images",
        ],
    ),
    (
        "Q&A",
        [],
    ),
]


def add_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    if bullets:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, line in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = BODY_COLOR
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left, top, width, height = Inches(1), Inches(2.5), Inches(11), Inches(2)
        box = slide.shapes.add_textbox(left, top, width, height)
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = ACCENT

    if slide.shapes.title is not None:
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = ACCENT
                run.font.size = Pt(32)
                run.font.bold = True


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for title, bullets in SLIDES:
        add_slide(prs, title, bullets)
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
